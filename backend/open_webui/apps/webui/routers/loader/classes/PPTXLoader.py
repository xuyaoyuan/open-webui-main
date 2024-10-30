from open_webui.apps.webui.routers.loader.classes.BaseLoader import BaseLoader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from spacy.lang.en import English
import os
import datetime

from spacy.lang.en import English
import shutil  
from typing import Dict, List, Union
import traceback
import re
from open_webui.apps.webui.routers.loader.classes.Chunk import Chunk
from open_webui.apps.webui.routers.loader.classes.Chunk import generate_unique_uuid

class PPTXLoader(BaseLoader):
    def __init__(self,
                 directory: str,
                 only_latest_content: bool = True,
                 recursive: bool = True,
                 encoding: str = 'utf-8',
                 exclude_files: list[str] = None,
                 num_threads: int = 1,
                 aiil: bool = False,
                 hitl: bool = False,
                 verbose: bool = False) -> None:
        
        super().__init__(directory=directory,
                         recursive=recursive, 
                         extensions=['.ppt', '.pptx'], 
                         encoding=encoding, 
                         exclude_files=exclude_files, 
                         num_threads=num_threads,
                         aiil=aiil,
                         hitl=hitl,
                         verbose=verbose)

        self.nlp = English()
        self.nlp.add_pipe("sentencizer")
        
        if not os.path.isdir(directory):
            raise Exception(f"`directory` must be a path of a directory, got {self.directory}")

    def _load_file(self, file_path: str):
        try:
            chunks = self._pptx_to_chunks(file_path)
            self.loaded_files[file_path] = chunks
        except Exception as e:
            print(f"\n***** Error reading {file_path}: {e} *****\n")

    def _pptx_to_chunks(self, file_path: str):
        try:
            document_id = ''
            match = re.search(r'attahments/(\d+)/', file_path)
            if match:
                document_id = match.group(1)

            prs = Presentation(file_path)
            chunks = []
            pptx_file_name = os.path.basename(file_path)
            slide_number = 0

            # Extract custom metadata
            core_properties = prs.core_properties
            author = core_properties.author or ''
            created_date = core_properties.created
            if created_date:
                created_date = created_date.strftime("%Y-%m-%d %H:%M:%S")
            else:
                created_date = ''

            for slide in prs.slides:
                slide_number += 1
                slide_contents = []
                slide_titles = []

                # Process each shape in the slide
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if shape.placeholder_format.type == 1 and shape.has_text_frame:
                            # Title placeholder
                            text = shape.text.strip()
                            if text:
                                slide_titles.append(text)
                                slide_contents.append(f"Title: {text}")
                    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        if shape.has_text_frame:
                            text = shape.text.strip()
                            if text:
                                # Use spaCy to split text into sentences
                                doc = self.nlp(text)
                                sentences = [sent.text.strip() for sent in doc.sents]
                                result_text = '\n\n'.join(sentences)
                                slide_contents.append(result_text)
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        image_bytes = image.blob
                        image_filename = f"{pptx_file_name}_slide{slide_number}_{image.filename}"
                        # Save the image to a directory (optional)
                        image_dir = os.path.dirname(file_path)  # Get the directory of the DOCX file
                        image_path = os.path.join(image_dir, image_filename)
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        slide_contents.append(f"[Image: {image_filename}]")
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            slide_contents.append("Table:\n" + table_text)
                    else:
                        # Handle other shape types if needed
                        pass

                # Extract notes (if any)
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_contents.append("Notes:\n" + notes_text)

                # Prepare metadata
                whole_titles = ', '.join(slide_titles)
                slide_meta = {
                    'pptx_file_name': pptx_file_name,
                    'slide_number': f'{slide_number}',
                    'page_title': whole_titles,
                    'author': author,
                    'created_date': created_date,
                    'source': pptx_file_name,
                    'confluence_id' : document_id
                }

                # Create Chunk
                chunk = Chunk(id=generate_unique_uuid())
                chunk.content = '\n'.join(slide_contents)
                chunk.metadata = self._filter_metadata(slide_meta)
                chunks.append(chunk)

            return chunks

        except Exception as e:
            import traceback
            traceback.print_exc()
            print(f"\n***** Error processing {file_path}: {e} *****\n")
            return None

    def _extract_table_text(self, table):
        table_text = ''
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_text.append(cell_text)
            if row_text:
                table_text += '\t'.join(row_text) + '\n'
        return table_text.strip()

    def _filter_metadata(self, metadata_dict: Dict[str, Union[str, List[str]]]) -> Dict[str, str]:
        """Convert list-type metadata values to comma-separated strings.

        Args:
            metadata_dict (Dict[str, Union[str, List[str]]]): Metadata.

        Returns:
            Dict[str, str]: Filtered metadata.
        """
        for key, value in metadata_dict.items():
            if isinstance(value, list):
                metadata_dict[key] = ', '.join(value)
        return metadata_dict
